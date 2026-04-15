import argparse
import json
import math
import re
import urllib.error
import urllib.request
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)
OLLAMA_MODEL = "llama3.2:3b"


DATASET_FALLBACKS = {
    "wikipedia_en_20220301": {
        "repo_id": "wikipedia",
        "notes": "English Wikipedia parquet snapshot from the legacy mirrored dataset repo.",
    },
    "wikipedia_en_20231101": {
        "repo_id": "wikimedia/wikipedia",
        "notes": "Current English Wikipedia parquet snapshot from the Wikimedia dataset repo.",
    },
    "open_web_math": {
        "repo_id": "open-web-math/open-web-math",
        "notes": "Large public math-heavy pretraining corpus.",
    },
    "codeparrot_clean": {
        "repo_id": "codeparrot/codeparrot-clean",
        "notes": "Deduplicated public Python code corpus.",
    },
    "github_code": {
        "repo_id": "codeparrot/github-code",
        "notes": "Multilanguage GitHub code dataset.",
    },
}


TEMPLATE_DEFINITIONS = [
    {"id": "pi-docs", "name": "Pi Docs", "tagline": "Blue editorial workspace", "font": "Aptos", "bg": (231, 241, 255), "primary": (13, 88, 187), "accent": (63, 155, 255), "accent_soft": (198, 226, 255), "text_muted": (56, 78, 110), "surface": (248, 252, 255)},
    {"id": "pi-sunset", "name": "Pi Sunset", "tagline": "Editorial and warm", "font": "Aptos", "bg": (248, 244, 236), "primary": (16, 29, 48), "accent": (197, 103, 55), "accent_soft": (232, 193, 171), "text_muted": (87, 96, 110), "surface": (255, 255, 255)},
    {"id": "midnight-grid", "name": "Midnight Grid", "tagline": "Sharp and technical", "font": "Aptos Display", "bg": (16, 20, 32), "primary": (240, 244, 255), "accent": (93, 196, 255), "accent_soft": (31, 44, 72), "text_muted": (164, 177, 204), "surface": (24, 31, 47)},
    {"id": "copper-paper", "name": "Copper Paper", "tagline": "Print-like narrative", "font": "Georgia", "bg": (245, 238, 229), "primary": (56, 38, 32), "accent": (174, 94, 63), "accent_soft": (225, 197, 176), "text_muted": (103, 86, 80), "surface": (255, 251, 247)},
    {"id": "ocean-glass", "name": "Ocean Glass", "tagline": "Cool and airy", "font": "Trebuchet MS", "bg": (236, 248, 249), "primary": (17, 70, 89), "accent": (0, 154, 173), "accent_soft": (182, 227, 233), "text_muted": (84, 108, 120), "surface": (255, 255, 255)},
    {"id": "citrus-pop", "name": "Citrus Pop", "tagline": "Bright and energetic", "font": "Verdana", "bg": (255, 247, 221), "primary": (61, 54, 11), "accent": (247, 148, 29), "accent_soft": (255, 216, 154), "text_muted": (110, 96, 32), "surface": (255, 252, 241)},
    {"id": "forest-note", "name": "Forest Note", "tagline": "Calm and grounded", "font": "Calibri", "bg": (239, 245, 239), "primary": (33, 61, 44), "accent": (90, 141, 93), "accent_soft": (197, 219, 198), "text_muted": (88, 106, 92), "surface": (255, 255, 255)},
    {"id": "rose-studio", "name": "Rose Studio", "tagline": "Soft and premium", "font": "Cambria", "bg": (252, 241, 242), "primary": (75, 33, 45), "accent": (212, 103, 130), "accent_soft": (241, 203, 213), "text_muted": (118, 83, 93), "surface": (255, 255, 255)},
    {"id": "terminal-neon", "name": "Terminal Neon", "tagline": "Command-center feel", "font": "Consolas", "bg": (10, 13, 20), "primary": (219, 255, 237), "accent": (74, 246, 169), "accent_soft": (21, 39, 31), "text_muted": (128, 182, 154), "surface": (17, 26, 23)},
    {"id": "lavender-ledger", "name": "Lavender Ledger", "tagline": "Soft analytics", "font": "Book Antiqua", "bg": (244, 240, 251), "primary": (59, 48, 95), "accent": (134, 98, 198), "accent_soft": (219, 208, 241), "text_muted": (101, 91, 131), "surface": (255, 255, 255)},
    {"id": "ruby-report", "name": "Ruby Report", "tagline": "Bold investor deck", "font": "Tahoma", "bg": (249, 238, 239), "primary": (76, 17, 30), "accent": (193, 47, 79), "accent_soft": (238, 189, 200), "text_muted": (110, 77, 84), "surface": (255, 255, 255)},
    {"id": "sandstone", "name": "Sandstone", "tagline": "Minimal and earthy", "font": "Segoe UI", "bg": (243, 238, 231), "primary": (74, 58, 39), "accent": (162, 121, 78), "accent_soft": (220, 203, 182), "text_muted": (113, 99, 82), "surface": (255, 252, 247)},
    {"id": "aurora", "name": "Aurora", "tagline": "Luminous and modern", "font": "Gill Sans MT", "bg": (235, 246, 248), "primary": (18, 52, 72), "accent": (122, 93, 255), "accent_soft": (202, 215, 255), "text_muted": (83, 101, 118), "surface": (255, 255, 255)},
    {"id": "graphite", "name": "Graphite", "tagline": "Executive monochrome", "font": "Century Gothic", "bg": (236, 237, 239), "primary": (31, 36, 44), "accent": (89, 99, 116), "accent_soft": (201, 207, 217), "text_muted": (94, 101, 111), "surface": (255, 255, 255)},
    {"id": "mint-board", "name": "Mint Board", "tagline": "Clean startup pitch", "font": "Candara", "bg": (237, 250, 245), "primary": (16, 64, 52), "accent": (39, 179, 132), "accent_soft": (187, 235, 218), "text_muted": (82, 112, 103), "surface": (255, 255, 255)},
    {"id": "solar-flare", "name": "Solar Flare", "tagline": "High-contrast keynote", "font": "Franklin Gothic Medium", "bg": (28, 24, 18), "primary": (255, 245, 225), "accent": (255, 133, 64), "accent_soft": (72, 49, 30), "text_muted": (209, 178, 145), "surface": (42, 34, 26)},
    {"id": "blueprint", "name": "Blueprint", "tagline": "Structured and precise", "font": "Corbel", "bg": (235, 242, 250), "primary": (27, 65, 112), "accent": (53, 131, 220), "accent_soft": (192, 218, 245), "text_muted": (84, 107, 132), "surface": (255, 255, 255)},
    {"id": "plum-theory", "name": "Plum Theory", "tagline": "Academic and rich", "font": "Palatino Linotype", "bg": (245, 239, 246), "primary": (70, 36, 78), "accent": (160, 86, 166), "accent_soft": (225, 197, 227), "text_muted": (112, 88, 118), "surface": (255, 255, 255)},
    {"id": "ember-slate", "name": "Ember Slate", "tagline": "Dark editorial", "font": "Bahnschrift", "bg": (28, 30, 37), "primary": (245, 242, 236), "accent": (224, 112, 73), "accent_soft": (74, 53, 46), "text_muted": (184, 170, 160), "surface": (41, 44, 54)},
    {"id": "skyline", "name": "Skyline", "tagline": "Corporate clean", "font": "Leelawadee UI", "bg": (242, 247, 251), "primary": (27, 57, 86), "accent": (0, 143, 213), "accent_soft": (196, 229, 247), "text_muted": (85, 103, 120), "surface": (255, 255, 255)},
    {"id": "studio-ink", "name": "Studio Ink", "tagline": "Magazine-inspired", "font": "Constantia", "bg": (247, 244, 238), "primary": (29, 29, 35), "accent": (221, 90, 43), "accent_soft": (239, 203, 187), "text_muted": (95, 92, 98), "surface": (255, 255, 255)},
]


def rgb(values: tuple[int, int, int]) -> RGBColor:
    return RGBColor(*values)


def get_template(template_id: str) -> dict:
    for template in TEMPLATE_DEFINITIONS:
        if template["id"] == template_id:
            return template
    return TEMPLATE_DEFINITIONS[0]


def slugify(value: str) -> str:
    cleaned = re.sub(r"[^a-zA-Z0-9]+", "-", value.strip().lower()).strip("-")
    return cleaned or "presentation"


def normalize_whitespace(value: str) -> str:
    return re.sub(r"\s+", " ", value).strip()


def split_sentences(value: str) -> list[str]:
    text = normalize_whitespace(value)
    if not text:
        return []
    parts = re.split(r"(?<=[.!?])\s+", text)
    return [part.strip(" -") for part in parts if part.strip(" -")]


def chunk_list(items: list[str], size: int) -> list[list[str]]:
    return [items[index : index + size] for index in range(0, len(items), size)]


def fit_text(value: str, limit: int) -> str:
    text = normalize_whitespace(value)
    if len(text) <= limit:
        return text
    return text[: max(0, limit - 3)].rstrip() + "..."


def convert_inline_slide_text(raw_text: str) -> str:
    text = raw_text.strip()
    if not text:
        return ""

    marker_pattern = re.compile(r"(?i)(slide\s+\d+\s*:)")
    parts = marker_pattern.split(text)
    if len(parts) <= 1:
        return text

    blocks = []
    current_marker = None
    current_content = []

    for part in parts:
        piece = part.strip()
        if not piece:
            continue
        if marker_pattern.fullmatch(piece):
            if current_marker is not None:
                blocks.append(f"{current_marker} {' '.join(current_content).strip()}".strip())
            current_marker = piece
            current_content = []
        else:
            current_content.append(normalize_whitespace(piece))

    if current_marker is not None:
        blocks.append(f"{current_marker} {' '.join(current_content).strip()}".strip())

    return "\n---\n".join(blocks)


def prettify_converted_slide_text(raw_text: str) -> str:
    normalized = convert_inline_slide_text(raw_text)
    if not normalized.strip():
        return ""

    pretty_blocks = []
    for block in normalized.split("\n---\n"):
        block = block.strip()
        if not block:
            continue

        match = re.match(r"(?i)^slide\s+(\d+)\s*:\s*(.*)$", block, re.DOTALL)
        if not match:
            pretty_blocks.append(block)
            continue

        slide_number = match.group(1)
        content = normalize_whitespace(match.group(2))
        sentences = split_sentences(content)
        if not sentences:
            pretty_blocks.append(f"Slide {slide_number}: Untitled")
            continue

        known_labels = [
            "Name",
            "Roll No",
            "Class/Batch",
            "Project Title",
            "Hardware",
            "Software",
            "Frontend",
            "Backend",
            "Scope",
            "Expected",
            "Actual",
        ]
        field_pattern = re.compile(
            r"(?i)(" + "|".join(re.escape(label) for label in known_labels) + r")\s*:\s*(.+?)(?=\s(?:"
            + "|".join(re.escape(label) for label in known_labels)
            + r")\s*:|$)"
        )
        field_matches = field_pattern.findall(content)
        field_bullets = []
        for key, value in field_matches[:4]:
            key_text = normalize_whitespace(key)
            value_text = normalize_whitespace(value)
            if key_text and value_text:
                field_bullets.append(f"{key_text}: {value_text}")

        title_source = content
        if field_matches:
            first_field_start = field_pattern.search(content)
            if first_field_start:
                title_source = content[: first_field_start.start()].strip()

        first_sentence = split_sentences(title_source)[0] if split_sentences(title_source) else title_source
        title = fit_text(first_sentence or "Untitled", 70)

        body_candidates = []
        if len(sentences) > 1:
            body_candidates.extend(sentences[1:3])
        body = fit_text(" ".join(body_candidates), 240)

        bullets = field_bullets[:3]
        if not bullets:
            bullets = [fit_text(item, 110) for item in sentences[1:4]]

        lines = [f"Slide {slide_number}: {title}"]
        if body:
            lines.append(body)
        for bullet in bullets[:3]:
            lines.append(f"- {fit_text(bullet, 110)}")
        pretty_blocks.append("\n".join(lines))

    return "\n---\n".join(pretty_blocks)


def call_ollama_json(prompt: str, model: str = OLLAMA_MODEL, timeout: int = 90) -> dict | None:
    payload = json.dumps(
        {
            "model": model,
            "prompt": prompt,
            "stream": False,
            "format": "json",
        }
    ).encode("utf-8")
    request = urllib.request.Request(
        "http://127.0.0.1:11434/api/generate",
        data=payload,
        headers={"Content-Type": "application/json"},
        method="POST",
    )
    try:
        with urllib.request.urlopen(request, timeout=timeout) as response:
            data = json.loads(response.read().decode("utf-8"))
    except (urllib.error.URLError, TimeoutError, json.JSONDecodeError, OSError):
        return None

    text = data.get("response", "").strip()
    if not text:
        return None
    try:
        return json.loads(text)
    except json.JSONDecodeError:
        return None


def shorten_story_title(title: str) -> str:
    text = normalize_whitespace(title)
    if len(text) <= 42:
        return text
    for separator in [" - ", ": ", ". "]:
        if separator in text:
            candidate = text.split(separator, 1)[0].strip()
            if 8 <= len(candidate) <= 42:
                return candidate
    words = text.split()
    return fit_text(" ".join(words[:6]), 42)


def enhance_records_with_llm(records: list[dict], model: str = OLLAMA_MODEL) -> list[dict]:
    enhanced = []
    for record in records:
        prompt = (
            "You are rewriting presentation slides.\n"
            "Return strict JSON with keys: title, subtitle, body, bullets, cue.\n"
            "Rules:\n"
            "- title: max 6 words\n"
            "- subtitle: max 10 words\n"
            "- body: max 28 words\n"
            "- bullets: array of up to 3 bullets, each max 10 words\n"
            "- cue: max 22 words\n"
            "- Do not output placeholder bullets.\n"
            "- Preserve concrete technical details when useful.\n"
            "- Avoid generic labels like Key Features unless truly necessary.\n"
            "- Keep the meaning but make it short and presentation-ready.\n\n"
            f"Input title: {record.get('title', '')}\n"
            f"Input subtitle: {record.get('subtitle', '')}\n"
            f"Input body: {record.get('body', '')}\n"
            f"Input bullets: {json.dumps(record.get('stats', []))}\n"
            f"Input cue: {record.get('cue', '')}\n"
        )
        response = call_ollama_json(prompt, model=model)
        if not response:
            enhanced.append(
                {
                    **record,
                    "title": shorten_story_title(record.get("title", "")),
                    "subtitle": fit_text(record.get("subtitle", ""), 60),
                    "body": fit_text(record.get("body", ""), 180),
                    "stats": [fit_text(item, 60) for item in record.get("stats", [])[:3]],
                    "cue": fit_text(record.get("cue", ""), 100),
                }
            )
            continue

        bullets = response.get("bullets") or record.get("stats", [])
        if not isinstance(bullets, list):
            bullets = record.get("stats", [])

        enhanced.append(
            {
                **record,
                "title": fit_text(response.get("title") or shorten_story_title(record.get("title", "")), 42),
                "subtitle": fit_text(response.get("subtitle") or record.get("subtitle", ""), 60),
                "body": fit_text(response.get("body") or record.get("body", ""), 180),
                "stats": [
                    fit_text(str(item), 60)
                    for item in bullets[:3]
                    if normalize_whitespace(str(item)) not in {"", "-", "--", "..."}
                ],
                "cue": fit_text(response.get("cue") or record.get("cue", ""), 100),
            }
        )
    return enhanced


def tidy_records(records: list[dict]) -> list[dict]:
    cleaned = []
    for record in records:
        title = shorten_story_title(record.get("title", ""))
        subtitle = fit_text(record.get("subtitle", ""), 60)
        body = fit_text(record.get("body", ""), 180)
        bullets = []
        for item in record.get("stats", [])[:3]:
            text = fit_text(str(item), 60)
            normalized = normalize_whitespace(text)
            if normalized and normalized not in {"-", "--", "..."}:
                bullets.append(text)
        cue = fit_text(record.get("cue", ""), 100)
        cleaned.append(
            {
                **record,
                "title": title,
                "subtitle": subtitle,
                "body": body,
                "stats": bullets,
                "cue": cue,
            }
        )
    return cleaned


def scan_presentation_records(records: list[dict], title: str = "") -> list[dict]:
    issues: list[dict] = []

    if len(title) > 60:
        issues.append(
            {
                "scope": "Deck",
                "severity": "warning",
                "message": "Presentation title is long and may look crowded on the cover slide.",
                "detail": fit_text(title, 90),
            }
        )

    if len(records) > 8:
        issues.append(
            {
                "scope": "Story Map",
                "severity": "error",
                "message": "Story Map will likely overflow because it lists more than 8 sections.",
                "detail": f"Current sections: {len(records)}",
            }
        )

    for index, record in enumerate(records, start=1):
        slide_label = f"Slide {index}"
        title_text = record.get("title", "")
        body_text = record.get("body", "")
        bullets = record.get("stats", [])
        cue_text = record.get("cue", "")

        if len(title_text) > 42:
            issues.append(
                {
                    "scope": slide_label,
                    "severity": "warning",
                    "message": "Title is longer than the recommended layout limit.",
                    "detail": fit_text(title_text, 90),
                }
            )

        if len(body_text) > 180:
            issues.append(
                {
                    "scope": slide_label,
                    "severity": "warning",
                    "message": "Body text is likely too dense for the content area.",
                    "detail": fit_text(body_text, 90),
                }
            )

        if len(bullets) > 3:
            issues.append(
                {
                    "scope": slide_label,
                    "severity": "warning",
                    "message": "More than 3 bullets reduces readability in the current layout.",
                    "detail": f"Bullet count: {len(bullets)}",
                }
            )

        for bullet in bullets[:5]:
            normalized_bullet = normalize_whitespace(str(bullet))
            if len(normalized_bullet) > 60:
                issues.append(
                    {
                        "scope": slide_label,
                        "severity": "warning",
                        "message": "A bullet point is too long and may wrap awkwardly.",
                        "detail": fit_text(normalized_bullet, 90),
                    }
                )
            if normalized_bullet in {"--", "-", "..."}:
                issues.append(
                    {
                        "scope": slide_label,
                        "severity": "error",
                        "message": "A bullet looks like placeholder or malformed text.",
                        "detail": normalized_bullet,
                    }
                )

        if len(cue_text) > 100:
            issues.append(
                {
                    "scope": slide_label,
                    "severity": "warning",
                    "message": "Speaker cue is longer than the cue card can comfortably hold.",
                    "detail": fit_text(cue_text, 90),
                }
            )

        dense_score = len(title_text) + len(body_text) + sum(len(str(item)) for item in bullets[:3])
        if dense_score > 260:
            issues.append(
                {
                    "scope": slide_label,
                    "severity": "error",
                    "message": "This slide is likely overcrowded for the chosen template.",
                    "detail": f"Density score: {dense_score}",
                }
            )

    return issues


def format_bytes(num_bytes: int) -> str:
    if num_bytes <= 0:
        return "0 B"
    units = ["B", "KB", "MB", "GB", "TB"]
    power = min(int(math.log(num_bytes, 1024)), len(units) - 1)
    value = num_bytes / (1024**power)
    return f"{value:.1f} {units[power]}"


def load_manifest(manifest_path: Path) -> dict:
    if not manifest_path.exists():
        return {"selected": [], "datasets": {}}
    return json.loads(manifest_path.read_text(encoding="utf-8"))


def collect_dataset_stats(dataset_dir: Path) -> dict:
    if not dataset_dir.exists():
        return {"exists": False, "files": 0, "size_bytes": 0, "top_types": []}

    files = [path for path in dataset_dir.rglob("*") if path.is_file()]
    suffix_counts: dict[str, int] = {}
    total_bytes = 0
    for path in files:
        suffix = path.suffix.lower() or "[no extension]"
        suffix_counts[suffix] = suffix_counts.get(suffix, 0) + 1
        total_bytes += path.stat().st_size

    top_types = sorted(suffix_counts.items(), key=lambda item: (-item[1], item[0]))[:3]
    return {"exists": True, "files": len(files), "size_bytes": total_bytes, "top_types": top_types}


def build_dataset_records(manifest: dict, output_dir: Path) -> list[dict]:
    dataset_specs = dict(manifest.get("datasets", {}))
    selected = manifest.get("selected") or list(dataset_specs.keys()) or list(DATASET_FALLBACKS.keys())
    records = []
    for name in selected:
        spec = dataset_specs.get(name, DATASET_FALLBACKS.get(name, {}))
        stats = collect_dataset_stats(output_dir / name)
        records.append(
            {
                "title": name.replace("_", " ").title(),
                "subtitle": spec.get("repo_id", "Unknown"),
                "body": spec.get("notes", "No dataset note available."),
                "stats": [
                    f"Files: {stats['files']:,}",
                    f"Disk usage: {format_bytes(stats['size_bytes'])}",
                    "Top types: " + (", ".join(f"{suffix} x{count}" for suffix, count in stats["top_types"]) or "Unavailable"),
                ],
                "cue": f"Present {name.replace('_', ' ')} as part of your training-data story. Lead with provenance, then explain why it matters.",
            }
        )
    return records


def parse_custom_slides(raw_text: str) -> list[dict]:
    raw_text = raw_text.strip()
    if not raw_text:
        return []

    slide_marker_pattern = re.compile(r"(?i)(?=slide\s+\d+\s*:)")
    if slide_marker_pattern.search(raw_text):
        chunks = [
            chunk.strip()
            for chunk in slide_marker_pattern.split(raw_text)
            if chunk.strip() and re.match(r"(?i)^slide\s+\d+\s*:", chunk.strip())
        ]
        slides = []
        for chunk in chunks:
            cleaned = re.sub(r"(?i)^slide\s+\d+\s*:\s*", "", chunk).strip()
            lines = [normalize_whitespace(line) for line in re.split(r"[\r\n]+", cleaned) if normalize_whitespace(line)]
            if len(lines) <= 1:
                sentences = split_sentences(cleaned)
                if not sentences:
                    continue
                title = fit_text(sentences[0], 70)
                subtitle = fit_text(sentences[1], 90) if len(sentences) > 1 else ""
                bullets = [fit_text(item, 110) for item in sentences[2:5]]
                body = fit_text(" ".join(sentences[1:3]), 240)
            else:
                title = fit_text(lines[0], 70)
                subtitle = fit_text(lines[1], 90) if len(lines) > 1 else ""
                remaining = lines[2:] if len(lines) > 2 else []
                bullet_candidates = []
                body_lines = []
                for item in remaining:
                    if item.startswith(("-", "*")):
                        bullet_candidates.append(item[1:].strip())
                    else:
                        body_lines.append(item)
                if not bullet_candidates:
                    bullet_candidates = split_sentences(" ".join(body_lines))[1:4]
                body = fit_text(" ".join(body_lines[:3]) or subtitle, 240)
                bullets = [fit_text(item, 110) for item in bullet_candidates[:3]]

            slides.append(
                {
                    "title": title,
                    "subtitle": subtitle,
                    "body": body,
                    "bullets": bullets,
                    "cue": fit_text(bullets[0] if bullets else body or title, 140),
                }
            )
        if slides:
            return slides

    if "\n---\n" not in raw_text and len(normalize_whitespace(raw_text)) > 320:
        sentences = split_sentences(raw_text)
        if len(sentences) > 4:
            slides = []
            first_sentence = sentences[0]
            sentence_chunks = chunk_list(sentences[1:], 4)
            for index, chunk in enumerate(sentence_chunks, start=1):
                body = " ".join(chunk[:2])
                bullets = chunk[2:] or chunk[:2]
                slides.append(
                    {
                        "title": fit_text(first_sentence if index == 1 else f"Key Points {index}", 70),
                        "subtitle": f"Auto-generated section {index}",
                        "body": fit_text(body, 240),
                        "bullets": [fit_text(item, 110) for item in bullets[:3]],
                        "cue": fit_text(chunk[0], 140) if chunk else "",
                    }
                )
            return slides

    slides = []
    blocks = [block.strip() for block in raw_text.split("\n---\n") if block.strip()]
    for block in blocks:
        lines = [line.rstrip() for line in block.splitlines() if line.strip()]
        if not lines:
            continue
        title = lines[0]
        subtitle = lines[1] if len(lines) > 1 and not lines[1].startswith(("-", "*")) else ""
        bullet_start = 2 if subtitle else 1
        bullets = []
        body_lines = []
        for line in lines[bullet_start:]:
            if line.startswith(("-", "*")):
                bullets.append(line[1:].strip())
            else:
                body_lines.append(line)
        if not bullets and body_lines:
            bullets = split_sentences(" ".join(body_lines))[1:4]
        slides.append(
            {
                "title": fit_text(title, 70),
                "subtitle": fit_text(subtitle, 90),
                "body": fit_text(" ".join(body_lines[:4]), 240),
                "bullets": [fit_text(item, 110) for item in bullets[:3]],
                "cue": fit_text(bullets[0] if bullets else "Use this slide to move the story forward.", 140),
            }
        )
    return slides


def build_custom_records(slides: list[dict]) -> list[dict]:
    records = []
    for slide in slides:
        title = slide.get("title", "").strip()
        if not title:
            continue
        bullets = [item.strip() for item in slide.get("bullets", []) if item.strip()]
        records.append(
            {
                "title": title,
                "subtitle": slide.get("subtitle", "").strip(),
                "body": slide.get("body", "").strip(),
                "stats": bullets[:3],
                "cue": slide.get("cue", "").strip() or "Use this slide to move the narrative forward with one clear takeaway.",
            }
        )
    return records


def build_presentation_payload(title: str, subtitle: str, records: list[dict], template_id: str, source_label: str) -> dict:
    return {"title": title, "subtitle": subtitle, "records": records, "template": get_template(template_id), "source_label": source_label}


def add_textbox(slide, left, top, width, height, text, font_size, color, font_name, bold=False, align=None):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.word_wrap = True
    paragraph = frame.paragraphs[0]
    paragraph.text = text
    paragraph.font.name = font_name
    paragraph.font.size = Pt(font_size)
    paragraph.font.bold = bold
    paragraph.font.color.rgb = color
    if align is not None:
        paragraph.alignment = align
    return box


def add_background(slide, template: dict) -> None:
    bg_fill = slide.background.fill
    bg_fill.solid()
    bg_fill.fore_color.rgb = rgb(template["bg"])

    if template["id"] == "pi-docs":
        frame = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.35), Inches(0.35), Inches(12.55), Inches(6.7))
        frame.fill.solid()
        frame.fill.fore_color.rgb = rgb(template["surface"])
        frame.line.color.rgb = rgb(template["accent"])

        left_pane = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.35), Inches(0.35), Inches(3.75), Inches(6.7))
        left_pane.fill.solid()
        left_pane.fill.fore_color.rgb = rgb(template["accent_soft"])
        left_pane.line.fill.background()

        overlay = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.35), Inches(0.35), Inches(3.75), Inches(6.7))
        overlay.fill.solid()
        overlay.fill.transparency = 0.28
        overlay.fill.fore_color.rgb = rgb(template["accent"])
        overlay.line.fill.background()

        divider = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(4.1), Inches(0.35), Inches(0.08), Inches(6.7))
        divider.fill.solid()
        divider.fill.fore_color.rgb = rgb(template["accent"])
        divider.line.fill.background()

        return

    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.6))
    band.fill.solid()
    band.fill.fore_color.rgb = rgb(template["primary"])
    band.line.fill.background()

    blob = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(11.55), Inches(-0.55), Inches(2.6), Inches(2.1))
    blob.fill.solid()
    blob.fill.fore_color.rgb = rgb(template["accent_soft"])
    blob.line.fill.background()

    stripe = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(7.05), SLIDE_WIDTH, Inches(0.45))
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = rgb(template["accent_soft"])
    stripe.line.fill.background()


def add_card(slide, left, top, width, height, template: dict):
    card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = rgb(template["surface"])
    card.line.color.rgb = rgb(template["accent_soft"])
    return card


def add_title_slide(prs: Presentation, payload: dict) -> None:
    template = payload["template"]
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, template)

    if template["id"] == "pi-docs":
        add_textbox(slide, Inches(0.78), Inches(0.78), Inches(2.5), Inches(0.35), "PI PRESENTATION", 11, rgb(template["surface"]), template["font"], bold=True)
        add_textbox(slide, Inches(0.78), Inches(4.8), Inches(2.5), Inches(0.35), "Created with DRMP", 12, rgb(template["surface"]), template["font"], bold=True)
        add_textbox(slide, Inches(0.78), Inches(5.18), Inches(2.8), Inches(0.5), payload["source_label"], 12, rgb(template["surface"]), template["font"])
        add_textbox(slide, Inches(4.7), Inches(2.65), Inches(6.2), Inches(2.2), payload["title"], 30, rgb(template["primary"]), template["font"], bold=True)
        slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(4.7), Inches(5.02), Inches(5.9), Inches(0.03)).fill.solid()
        line = slide.shapes[-1]
        line.fill.fore_color.rgb = rgb(template["primary"])
        line.line.fill.background()
        add_textbox(slide, Inches(4.7), Inches(5.28), Inches(4.5), Inches(0.4), "Author: DRMP User", 14, rgb(template["text_muted"]), template["font"], bold=True)
        add_textbox(slide, Inches(4.7), Inches(5.58), Inches(4.5), Inches(0.35), payload["subtitle"], 11, rgb(template["text_muted"]), template["font"])
        return

    badge = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(0.95), Inches(2.5), Inches(0.46))
    badge.fill.solid()
    badge.fill.fore_color.rgb = rgb(template["accent"])
    badge.line.fill.background()
    badge.text_frame.text = template["name"].upper()
    badge.text_frame.paragraphs[0].font.name = template["font"]
    badge.text_frame.paragraphs[0].font.size = Pt(13)
    badge.text_frame.paragraphs[0].font.bold = True
    badge.text_frame.paragraphs[0].font.color.rgb = rgb(template["surface"])
    badge.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    add_textbox(slide, Inches(0.8), Inches(1.65), Inches(8.3), Inches(1.35), payload["title"], 28, rgb(template["primary"]), template["font"], bold=True)
    add_textbox(slide, Inches(0.8), Inches(3.0), Inches(8.1), Inches(1.0), payload["subtitle"], 16, rgb(template["text_muted"]), template["font"])

    add_card(slide, Inches(9.2), Inches(1.25), Inches(3.1), Inches(4.95), template)
    add_textbox(slide, Inches(9.55), Inches(1.65), Inches(2.3), Inches(0.35), "Deck Profile", 14, rgb(template["accent"]), template["font"], bold=True)
    add_textbox(slide, Inches(9.55), Inches(2.15), Inches(2.25), Inches(1.8), f"{template['tagline']}\n20 selectable templates\nPPT export ready", 18, rgb(template["primary"]), template["font"], bold=True)
    add_textbox(slide, Inches(9.55), Inches(5.5), Inches(2.0), Inches(0.45), payload["source_label"], 11, rgb(template["text_muted"]), template["font"])


def add_summary_slide(prs: Presentation, payload: dict) -> None:
    template = payload["template"]
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, template)

    if template["id"] == "pi-docs":
        add_textbox(slide, Inches(0.75), Inches(0.78), Inches(4.8), Inches(0.45), "CONTENTS", 28, rgb(template["primary"]), template["font"], bold=True)
        top_line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.75), Inches(1.62), Inches(10.0), Inches(0.03))
        top_line.fill.solid()
        top_line.fill.fore_color.rgb = rgb(template["primary"])
        top_line.line.fill.background()

        summary_titles = [shorten_story_title(record["title"]) for record in payload["records"]]
        current_titles = summary_titles[:10]
        add_card(slide, Inches(0.75), Inches(1.9), Inches(10.05), Inches(3.1), template)
        left_titles = current_titles[::2]
        right_titles = current_titles[1::2]
        y_left = 2.15
        for idx, record_title in enumerate(left_titles, start=1):
            number = idx * 2 - 1
            add_card(slide, Inches(0.95), Inches(y_left), Inches(0.42), Inches(0.32), template)
            add_textbox(slide, Inches(1.07), Inches(y_left + 0.03), Inches(0.2), Inches(0.2), str(number), 16, rgb(template["primary"]), template["font"], bold=True)
            add_textbox(slide, Inches(1.52), Inches(y_left), Inches(3.7), Inches(0.32), record_title, 16, rgb(template["primary"]), template["font"], bold=True)
            y_left += 0.62
        y_right = 2.15
        for idx, record_title in enumerate(right_titles, start=1):
            number = idx * 2
            add_card(slide, Inches(5.78), Inches(y_right), Inches(0.42), Inches(0.32), template)
            add_textbox(slide, Inches(5.9), Inches(y_right + 0.03), Inches(0.2), Inches(0.2), str(number), 16, rgb(template["primary"]), template["font"], bold=True)
            add_textbox(slide, Inches(6.35), Inches(y_right), Inches(3.9), Inches(0.32), record_title, 16, rgb(template["primary"]), template["font"], bold=True)
            y_right += 0.62
        return

    add_textbox(slide, Inches(0.8), Inches(0.95), Inches(5.8), Inches(0.6), "Story Map", 24, rgb(template["primary"]), template["font"], bold=True)

    metrics = [("Slides", str(len(payload["records"]) + 3)), ("Sections", str(len(payload["records"]))), ("Theme", template["name"])]
    left = 0.8
    for label, value in metrics:
        add_card(slide, Inches(left), Inches(1.8), Inches(2.85), Inches(1.55), template)
        add_textbox(slide, Inches(left + 0.22), Inches(2.03), Inches(2.2), Inches(0.3), label, 12, rgb(template["text_muted"]), template["font"])
        add_textbox(slide, Inches(left + 0.22), Inches(2.38), Inches(2.35), Inches(0.52), value, 22, rgb(template["primary"]), template["font"], bold=True)
        left += 3.05

    summary_titles = [shorten_story_title(record["title"]) for record in payload["records"]]
    summary_pages = chunk_list(summary_titles, 4)
    current_titles = summary_pages[0] if summary_pages else []

    add_card(slide, Inches(0.8), Inches(3.75), Inches(11.45), Inches(2.45), template)
    add_textbox(slide, Inches(1.08), Inches(4.03), Inches(3.2), Inches(0.35), "Included narrative beats", 14, rgb(template["accent"]), template["font"], bold=True)
    y = 4.45
    for record_title in current_titles:
        add_textbox(slide, Inches(1.08), Inches(y), Inches(10.2), Inches(0.32), f"- {record_title}", 16, rgb(template["primary"]), template["font"])
        y += 0.42

    for page_index, page_titles in enumerate(summary_pages[1:], start=2):
        extra_slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_background(extra_slide, template)
        add_textbox(extra_slide, Inches(0.8), Inches(0.95), Inches(7.0), Inches(0.6), f"Story Map Continued", 24, rgb(template["primary"]), template["font"], bold=True)
        add_card(extra_slide, Inches(0.8), Inches(1.9), Inches(11.45), Inches(3.8), template)
        add_textbox(extra_slide, Inches(1.08), Inches(2.18), Inches(5.5), Inches(0.35), f"Included narrative beats · Page {page_index}", 14, rgb(template["accent"]), template["font"], bold=True)
        y = 2.65
        for record_title in page_titles:
            add_textbox(extra_slide, Inches(1.08), Inches(y), Inches(10.0), Inches(0.36), f"- {record_title}", 18, rgb(template["primary"]), template["font"])
            y += 0.58


def add_content_slide(prs: Presentation, record: dict, template: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, template)

    if template["id"] == "pi-docs":
        add_textbox(slide, Inches(0.85), Inches(2.2), Inches(5.8), Inches(0.8), record["title"], 28, rgb(template["primary"]), template["font"], bold=True)
        add_textbox(slide, Inches(0.85), Inches(3.2), Inches(5.7), Inches(2.0), record.get("body", ""), 15, rgb(template["text_muted"]), template["font"])
        if record.get("stats"):
            add_card(slide, Inches(0.85), Inches(5.1), Inches(2.8), Inches(1.2), template)
            add_textbox(slide, Inches(1.05), Inches(5.32), Inches(2.2), Inches(0.25), "Key Points", 11, rgb(template["accent"]), template["font"], bold=True)
            y = 5.56
            for bullet in record.get("stats", [])[:2]:
                add_textbox(slide, Inches(1.05), Inches(y), Inches(2.2), Inches(0.22), f"- {bullet}", 10, rgb(template["text_muted"]), template["font"])
                y += 0.22
        cue_card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(8.0), Inches(0.7), Inches(3.8), Inches(5.4))
        cue_card.fill.solid()
        cue_card.fill.fore_color.rgb = rgb(template["accent_soft"])
        cue_card.line.fill.background()
        shade = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(8.0), Inches(0.7), Inches(3.8), Inches(5.4))
        shade.fill.solid()
        shade.fill.transparency = 0.35
        shade.fill.fore_color.rgb = rgb(template["primary"])
        shade.line.fill.background()
        add_textbox(slide, Inches(8.3), Inches(4.95), Inches(2.7), Inches(0.25), "Commentary", 11, rgb(template["surface"]), template["font"], bold=True)
        add_textbox(slide, Inches(8.3), Inches(5.22), Inches(2.9), Inches(0.5), record.get("cue", ""), 12, rgb(template["surface"]), template["font"])
        return

    title_size = 24 if len(record["title"]) < 48 else 20
    body_size = 17 if len(record.get("body", "")) < 170 else 14
    cue_size = 18 if len(record.get("cue", "")) < 110 else 15

    add_textbox(slide, Inches(0.8), Inches(0.95), Inches(8.6), Inches(0.8), record["title"], title_size, rgb(template["primary"]), template["font"], bold=True)
    add_textbox(slide, Inches(0.8), Inches(1.62), Inches(8.1), Inches(0.32), record.get("subtitle", ""), 12, rgb(template["accent"]), template["font"], bold=True)
    add_textbox(slide, Inches(0.8), Inches(2.0), Inches(7.8), Inches(1.45), record.get("body", ""), body_size, rgb(template["text_muted"]), template["font"])

    add_card(slide, Inches(0.8), Inches(3.3), Inches(4.45), Inches(2.15), template)
    add_textbox(slide, Inches(1.08), Inches(3.58), Inches(3.3), Inches(0.35), "Key Points", 14, rgb(template["accent"]), template["font"], bold=True)
    y = 4.0
    for bullet in record.get("stats", [])[:3]:
        bullet_size = 14 if len(bullet) < 72 else 12
        add_textbox(slide, Inches(1.08), Inches(y), Inches(3.7), Inches(0.42), f"- {bullet}", bullet_size, rgb(template["primary"]), template["font"])
        y += 0.42

    cue_card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(5.7), Inches(3.0), Inches(5.85), Inches(2.8))
    cue_card.fill.solid()
    cue_card.fill.fore_color.rgb = rgb(template["primary"])
    cue_card.line.fill.background()
    add_textbox(slide, Inches(6.0), Inches(3.32), Inches(4.8), Inches(0.35), "Speaker Cue", 14, rgb(template["accent_soft"]), template["font"], bold=True)
    add_textbox(slide, Inches(6.0), Inches(3.78), Inches(4.95), Inches(1.75), record.get("cue", ""), cue_size, rgb(template["surface"]), template["font"], bold=True)


def add_closing_slide(prs: Presentation, payload: dict) -> None:
    template = payload["template"]
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, template)
    add_textbox(slide, Inches(0.8), Inches(1.2), Inches(8.0), Inches(0.6), "Export Ready", 26, rgb(template["primary"]), template["font"], bold=True)
    add_textbox(slide, Inches(0.8), Inches(2.05), Inches(9.4), Inches(2.2), "This deck was generated automatically from your selected template and slide content. Use it as-is, or keep iterating from the web UI to try other looks before sharing.", 18, rgb(template["text_muted"]), template["font"])
    add_textbox(slide, Inches(0.8), Inches(5.42), Inches(7.0), Inches(0.35), f"Theme: {template['name']} | Source: {payload['source_label']}", 12, rgb(template["accent"]), template["font"], bold=True)


def generate_presentation(payload: dict, pptx_path: Path) -> Path:
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    add_title_slide(prs, payload)
    add_summary_slide(prs, payload)
    for record in payload["records"]:
        add_content_slide(prs, record, payload["template"])
    add_closing_slide(prs, payload)
    pptx_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(pptx_path))
    return pptx_path


def create_dataset_payload(manifest_path: Path, output_dir: Path, template_id: str, title: str | None = None, subtitle: str | None = None) -> dict:
    manifest = load_manifest(manifest_path)
    records = build_dataset_records(manifest, output_dir)
    return build_presentation_payload(title or "Dataset Download Overview", subtitle or "A polished PowerPoint deck generated from local dataset folders and the saved manifest.", records, template_id, f"Dataset mode | {output_dir}")


def main() -> None:
    parser = argparse.ArgumentParser(description="Generate a PowerPoint presentation.")
    parser.add_argument("--manifest", default="datasets/download_manifest.json", help="Path to the dataset manifest JSON file.")
    parser.add_argument("--output-dir", default="datasets", help="Folder containing downloaded datasets.")
    parser.add_argument("--pptx", default="datasets/dataset_overview_presentation.pptx", help="Destination PowerPoint file.")
    parser.add_argument("--template", default="pi-sunset", choices=[item["id"] for item in TEMPLATE_DEFINITIONS], help="Template identifier.")
    parser.add_argument("--title", default=None, help="Optional presentation title override.")
    parser.add_argument("--subtitle", default=None, help="Optional presentation subtitle override.")
    parser.add_argument("--slides-text", default=None, help="Optional custom slides text separated by lines containing ---.")
    parser.add_argument("--use-llm", action="store_true", help="Use local Ollama to shorten and polish slide content.")
    args = parser.parse_args()

    if args.slides_text:
        records = build_custom_records(parse_custom_slides(args.slides_text))
        if args.use_llm:
            records = enhance_records_with_llm(records)
        payload = build_presentation_payload(args.title or "Custom Presentation", args.subtitle or "Generated from freeform slide content.", records, args.template, "Custom mode")
    else:
        payload = create_dataset_payload(Path(args.manifest), Path(args.output_dir), args.template, args.title, args.subtitle)

    generate_presentation(payload, Path(args.pptx))
    print(f"PowerPoint saved to {args.pptx}")


if __name__ == "__main__":
    main()
